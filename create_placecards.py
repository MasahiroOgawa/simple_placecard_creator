#!/usr/bin/env python3
"""Create wedding place cards (席札) as a LibreOffice-editable .pptx and PDF.

Reads names from an Excel (.xlsx), CSV, or plain text file and produces
A4-portrait pages with 8 cards each (2x4). Each card is designed to be
folded in half — text appears in the bottom half only.
"""

import argparse
import csv
import math
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Mm, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# --- Layout constants ---
A4_WIDTH = 210.0    # mm (portrait)
A4_HEIGHT = 297.0   # mm (portrait)
COLS = 2
ROWS = 4
CARDS_PER_PAGE = COLS * ROWS
CARD_WIDTH = A4_WIDTH / COLS      # 105mm
CARD_HEIGHT = A4_HEIGHT / ROWS    # 74.25mm

# --- Style defaults ---
FONT_NAME = "Noto Serif CJK JP"
SCRIPT_FONT_NAME = "Great Vibes"  # curly arabesque-style script font
WELCOME_TEXT = "welcome"
DATE_TEXT = "April 19, 2026"
WELCOME_FONT_SIZE = Pt(27)
DATE_FONT_SIZE = Pt(14)
NAME_FONT_SIZES = {6: Pt(28), 10: Pt(24)}  # max_len: size; fallback Pt(20)
NAME_COLOR = RGBColor(0x33, 0x33, 0x33)
SUB_COLOR = RGBColor(0x80, 0x80, 0x80)
LINE_COLOR = RGBColor(0xC0, 0xC0, 0xC0)
LINE_WIDTH = Pt(0.5)
TEXT_BLOCK_HEIGHT_MM = 38


def load_names_from_xlsx(xlsx_path):
    """Load attending guest names from xlsx (ゲスト一覧 format).

    Extracts main guest name and 連名 (joint names), filtering to 'ご出席' only.
    Appends 様 to each name.
    """
    from openpyxl import load_workbook
    wb = load_workbook(xlsx_path, read_only=True)
    ws = wb.active
    rows = list(ws.iter_rows(values_only=True))
    header = list(rows[0])
    names = []
    for row in rows[1:]:
        data = dict(zip(header, row))
        if data.get('出欠情報') != 'ご出席':
            continue
        guest = data.get('ゲスト名', '')
        if guest:
            names.append(f'{guest}様')
        # Add joint names (連名1-4)
        for i in range(1, 5):
            joint = data.get(f'連名{i}')
            if joint:
                names.append(f'{joint}様')
    wb.close()
    return names


def load_names_from_csv(csv_path):
    """Load names from CSV (expects '様付き' column or '表示名' column)."""
    names = []
    with open(csv_path, encoding='utf-8') as f:
        reader = csv.DictReader(f)
        for row in reader:
            if '様付き' in row:
                names.append(row['様付き'])
            elif '表示名' in row:
                names.append(f'{row["表示名"]}様')
    return names


def load_names_from_txt(txt_path):
    """Load names from plain text file (one name per line, 様 appended if missing)."""
    names = []
    with open(txt_path, encoding='utf-8') as f:
        for line in f:
            line = line.strip()
            if line:
                names.append(line if line.endswith('様') else f'{line}様')
    return names


def name_font_size(name):
    """Pick font size based on name length."""
    for max_len, size in sorted(NAME_FONT_SIZES.items()):
        if len(name) <= max_len:
            return size
    return Pt(20)


def add_separator_lines(slide):
    """Add gray cut lines between cards and dashed fold lines within each card."""
    # Vertical cut line between columns
    x = Mm(CARD_WIDTH)
    c = slide.shapes.add_connector(1, x, Mm(0), x, Mm(A4_HEIGHT))
    c.line.color.rgb = LINE_COLOR
    c.line.width = LINE_WIDTH
    # Horizontal cut lines between rows
    for row in range(1, ROWS):
        y = Mm(row * CARD_HEIGHT)
        c = slide.shapes.add_connector(1, Mm(0), y, Mm(A4_WIDTH), y)
        c.line.color.rgb = LINE_COLOR
        c.line.width = LINE_WIDTH
    # Dashed fold lines at the vertical center of each card
    for row in range(ROWS):
        y = Mm(row * CARD_HEIGHT + CARD_HEIGHT / 2)
        c = slide.shapes.add_connector(1, Mm(0), y, Mm(A4_WIDTH), y)
        c.line.color.rgb = LINE_COLOR
        c.line.width = LINE_WIDTH
        c.line.dash_style = 2  # dash


def create_placecards(names, output_pptx, welcome=WELCOME_TEXT, date=DATE_TEXT, font=FONT_NAME):
    """Generate the .pptx place-card file (8 cards per A4 portrait page, 2x4).

    Each card is folded in half — text is placed in the bottom half of each card.
    """
    prs = Presentation()
    prs.slide_width = Mm(A4_WIDTH)
    prs.slide_height = Mm(A4_HEIGHT)

    # Text is centered in the bottom half of each card
    half_card = CARD_HEIGHT / 2  # ~37.1mm
    text_top_in_card = half_card + (half_card - TEXT_BLOCK_HEIGHT_MM) / 2
    num_pages = math.ceil(len(names) / CARDS_PER_PAGE)

    for page in range(num_pages):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        add_separator_lines(slide)

        for slot in range(CARDS_PER_PAGE):
            idx = page * CARDS_PER_PAGE + slot
            if idx >= len(names):
                break

            name = names[idx]
            col = slot % COLS
            row = slot // COLS

            left = Mm(col * CARD_WIDTH)
            top = Mm(row * CARD_HEIGHT + text_top_in_card)
            width = Mm(CARD_WIDTH)
            height = Mm(TEXT_BLOCK_HEIGHT_MM)

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True

            # Welcome line
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.space_before = Pt(0)
            p.space_after = Pt(4)
            run = p.add_run()
            run.text = welcome
            run.font.name = SCRIPT_FONT_NAME
            run.font.size = WELCOME_FONT_SIZE
            run.font.color.rgb = SUB_COLOR

            # Name line
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            p.space_before = Pt(0)
            p.space_after = Pt(4)
            run = p.add_run()
            run.text = name
            run.font.name = font
            run.font.size = name_font_size(name)
            run.font.color.rgb = NAME_COLOR

            # Date line
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            p.space_before = Pt(0)
            p.space_after = Pt(0)
            run = p.add_run()
            run.text = date
            run.font.name = SCRIPT_FONT_NAME
            run.font.size = DATE_FONT_SIZE
            run.font.color.rgb = SUB_COLOR

    prs.save(str(output_pptx))
    return num_pages


def convert_to_pdf(pptx_path):
    """Convert .pptx to PDF using LibreOffice."""
    outdir = str(Path(pptx_path).parent)
    result = subprocess.run(
        ['libreoffice', '--headless', '--convert-to', 'pdf', str(pptx_path), '--outdir', outdir],
        capture_output=True, text=True,
    )
    pdf_path = Path(pptx_path).with_suffix('.pdf')
    if pdf_path.exists():
        return pdf_path
    print(f'Warning: PDF conversion may have failed.\n{result.stderr}', file=sys.stderr)
    return None


def main():
    parser = argparse.ArgumentParser(description='Create place cards (席札) from a name list')
    parser.add_argument('input', help='Excel (.xlsx), CSV, or plain text file with names')
    parser.add_argument('-o', '--output', default=None, help='Output directory (default: output/)')
    parser.add_argument('--welcome', default=WELCOME_TEXT, help=f'Welcome text (default: {WELCOME_TEXT})')
    parser.add_argument('--date', default=DATE_TEXT, help=f'Date text (default: {DATE_TEXT})')
    parser.add_argument('--font', default=FONT_NAME, help=f'Font name (default: {FONT_NAME})')
    parser.add_argument('--no-pdf', action='store_true', help='Skip PDF conversion')
    args = parser.parse_args()

    input_path = Path(args.input)
    if input_path.suffix == '.xlsx':
        names = load_names_from_xlsx(input_path)
    elif input_path.suffix == '.csv':
        names = load_names_from_csv(input_path)
    else:
        names = load_names_from_txt(input_path)

    if not names:
        print('No names found in input file.', file=sys.stderr)
        sys.exit(1)

    output_dir = Path(args.output) if args.output else Path('output')
    output_dir.mkdir(parents=True, exist_ok=True)
    output_pptx = output_dir / 'placecards.pptx'
    num_pages = create_placecards(names, output_pptx, welcome=args.welcome, date=args.date, font=args.font)
    print(f'Created {output_pptx} ({len(names)} names, {num_pages} pages)')

    if not args.no_pdf:
        pdf_path = convert_to_pdf(output_pptx)
        if pdf_path:
            print(f'Created {pdf_path}')


if __name__ == '__main__':
    main()
